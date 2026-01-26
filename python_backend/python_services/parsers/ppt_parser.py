"""
PPT解析器 - 将PPT转换为结构化MD文件
todo: 以后决定一下是转为md文件呢（图片可以通过提示词微调）？还是直接拿langchain的转(没有图片)
todo: 调整一下图片的内容与选择
"""
import os
import uuid
from pathlib import Path
from typing import List, Dict, Any
from datetime import datetime
from zoneinfo import ZoneInfo
from PIL import Image

from langchain_core.documents import Document
from python_services.parsers.base_parser import BaseParser
from python_services.utils.ocr_util import OcrUtil

# 尝试导入python-pptx库，如未安装则使用备用方案
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
    print("python-pptx库已加载，使用增强PPT解析模式")
except ImportError:
    print("python-pptx库未安装，使用基础解析模式")
    PPTX_AVAILABLE = False
    from langchain_community.document_loaders import UnstructuredPowerPointLoader


class PPtParser(BaseParser):
    """PPT解析器，将PPT转换为结构化MD文件"""

    def __init__(self, vision_llm=None):
        super().__init__("PPT解析器", ["pptx", "ppt"])
        self.vision_llm = vision_llm
        
    def parse_impl(self, file_path_or_url: str) -> list[Document]:
        """解析PPT的主要函数，转换为结构化MD文件"""
        file_path = Path(file_path_or_url)
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        documents = []
        file_stem = file_path.stem
        
        # 创建输出目录
        output_dir = Path("./output") / file_stem
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # 根据是否有python-pptx库选择解析方式
        if PPTX_AVAILABLE:
            # 使用增强解析模式
            md_content = self._parse_ppt_enhanced(file_path, output_dir)
        else:
            # 使用基础解析模式
            md_content = self._parse_ppt_basic(file_path)
        
        # 生成MD文件
        md_file_path = output_dir / f"{file_stem}.md"
        with open(md_file_path, "w", encoding="utf-8") as f:
            f.write(md_content)
        print(f"PPT已转换为MD文件: {md_file_path}")
        
        # 创建Document对象
        document = Document(
            page_content=md_content,
            metadata={
                "source": str(file_path),
                "type": "text",
                "file_name": file_path.name,
                "ext": "md",
                "create_time": datetime.now(ZoneInfo("Asia/Shanghai")).isoformat(),
                "total_slides": md_content.count("# 幻灯片")  # 估算幻灯片数量
            }
        )
        documents.append(document)
        
        return documents
    
    def _parse_ppt_enhanced(self, file_path: Path, output_dir: Path) -> str:
        """增强PPT解析，使用python-pptx库提取结构化信息"""
        prs = Presentation(str(file_path))
        md_parts = []
        
        # 提取PPT基本信息
        md_parts.append(f"# {file_path.stem}")
        # md_parts.append(f"> 转换时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        # md_parts.append(f"> 幻灯片数量: {len(prs.slides)}")
        # md_parts.append("")
        
        # 遍历每页幻灯片
        for slide_idx, slide in enumerate(prs.slides, 1):
            # md_parts.append(f"## 幻灯片 {slide_idx}")
            
            # 提取标题
            slide_title = self._get_slide_title(slide)
            if slide_title:
                md_parts.append(f"## {slide_title}")
            
            # 提取文本内容
            text_content = self._get_slide_text(slide)
            if text_content:
                md_parts.append(text_content)
            
            # 提取图片
            images_info = self._extract_slide_images(slide, file_path.stem, output_dir, slide_idx)
            if images_info:
                for img_info in images_info:
                    # 添加图片引用
                    md_parts.append(f"### 图片路径及内容\n![幻灯片{slide_idx}图片]({img_info['path']})")
                    
                    # 添加图片OCR结果
                    if img_info['ocr_result']:
                        # md_parts.append(f"\n### 图片内容识别")
                        md_parts.append(f"```")
                        md_parts.append(img_info['ocr_result'])
                        md_parts.append(f"```")
                        md_parts.append("")
                    else:
                        # md_parts.append(f"\n### 图片内容识别")
                        md_parts.append(f"> 警告：图片无法识别或未包含可提取文本")
                        md_parts.append("")
            
            md_parts.append("")  # 幻灯片之间添加空行
        
        return "\n".join(md_parts)
    
    def _parse_ppt_basic(self, file_path: Path) -> str:
        """基础PPT解析，使用UnstructuredPowerPointLoader"""
        loader = UnstructuredPowerPointLoader(str(file_path))
        loader_docs = loader.load()
        
        md_parts = [f"# {file_path.stem}"]
        md_parts.append(f"> 转换时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        md_parts.append("")
        
        # 将所有内容合并为MD格式
        for idx, doc in enumerate(loader_docs, 1):
            md_parts.append(f"## 幻灯片 {idx}")
            md_parts.append(doc.page_content)
            md_parts.append("")
        
        return "\n".join(md_parts)
    
    def _get_slide_title(self, slide) -> str:
        """获取幻灯片标题"""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame.text:
                # 通常标题是第一个文本框
                return shape.text_frame.text.strip()
        return ""
    
    def _get_slide_text(self, slide) -> str:
        """获取幻灯片文本内容"""
        text_parts = []
        title_processed = False
        
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            
            text_frame = shape.text_frame
            if not text_frame.text:
                continue
            
            # 跳过标题
            if not title_processed:
                title_processed = True
                continue
            
            # 处理段落
            for paragraph in text_frame.paragraphs:
                if not paragraph.text:
                    continue
                    
                # 处理列表
                if paragraph.level > 0:
                    text_parts.append(f"{'  ' * paragraph.level}- {paragraph.text}")
                else:
                    text_parts.append(paragraph.text)
        
        return "\n".join(text_parts)
    
    def _extract_slide_images(self, slide, file_stem: str, output_dir: Path, slide_idx: int) -> List[Dict[str, Any]]:
        """提取幻灯片中的图片并进行OCR识别
        
        :return: 包含图片路径和OCR结果的字典列表
        """
        images_info = []
        image_dir = output_dir / "images"
        image_dir.mkdir(exist_ok=True)
        
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    # 提取图片，添加异常处理以防止某些特殊图片导致崩溃
                    image_stream = shape.image.blob
                    image_ext = shape.image.ext
                    image_name = f"slide_{slide_idx}_image_{shape_idx}_{uuid.uuid4()}.{image_ext}"
                    image_path = image_dir / image_name
                    
                    # 保存图片
                    with open(image_path, "wb") as f:
                        f.write(image_stream)
                    
                    # 打开图片进行预处理和分析
                    with Image.open(image_path) as img:
                        # 获取图片尺寸
                        width, height = img.size
                        
                        # 过滤小图标（尺寸小于20x20像素）
                        if width < 20 and height < 20:
                            print(f"跳过小图标（尺寸：{width}x{height}）: {image_path}")
                            continue
                        
                        # 对图片进行OCR/VLM识别
                        ocr_result = self._perform_image_ocr(img)
                    
                    # 添加到图片信息列表
                    relative_path = f"images/{image_name}"
                    images_info.append({
                        'path': relative_path,
                        'absolute_path': image_path,
                        'width': width,
                        'height': height,
                        'ocr_result': ocr_result
                    })
                    print(f"已提取图片: {image_path}")
                    if ocr_result:
                        print(f"  OCR识别成功，提取文本长度: {len(ocr_result)}")
                except AttributeError as e:
                    # 处理'Part' object has no attribute 'image'错误
                    print(f"警告：跳过幻灯片{slide_idx}中无法提取的图片（索引{shape_idx}）: {e}")
                    continue
                except Exception as e:
                    # 处理其他可能的图片提取错误
                    print(f"警告：幻灯片{slide_idx}中图片（索引{shape_idx}）提取失败: {e}")
                    continue
        
        return images_info
    
    def _perform_image_ocr(self, image: Image.Image) -> str:
        """对图片进行OCR/VLM识别
        
        :param image: PIL.Image对象
        :return: 识别的文本结果
        """
        try:
            # 优先使用视觉LLM进行OCR
            ocr_result = OcrUtil.vision_ocr(self.vision_llm, image)
            
            # 如果视觉LLM不可用或识别结果为空，使用Tesseract OCR作为备选
            if not ocr_result:
                ocr_result = OcrUtil.tesseract_ocr(image)
            
            return ocr_result
        except ConnectionError as e:
            print(f"OCR识别失败：网络连接异常 - {e}")
            print("  解决方案建议：检查网络连接，确保服务可访问")
            return ""
        except ValueError as e:
            print(f"OCR识别失败：图片格式错误 - {e}")
            print("  解决方案建议：检查图片格式，确保为JPEG/PNG等支持格式")
            return ""
        except AttributeError as e:
            print(f"OCR识别失败：视觉模型接口错误 - {e}")
            print("  解决方案建议：检查视觉模型配置，确保API接口可用")
            return ""
        except Exception as e:
            print(f"OCR识别失败：未知错误 - {e}")
            print("  解决方案建议：检查图片质量，确保图片清晰可读")
            import traceback
            traceback.print_exc()
            return ""


if __name__ == '__main__':
    parser = PPtParser(vision_llm=None)
    ppt_path = r"C:\Users\ASUS\Desktop\中国石油大学（北京）本科招生宣传.pptx"
    documents = parser.parse(ppt_path)
    print(f"\n解析完成，生成 {len(documents)} 个文档")
    print(f"文档内容预览（前500字符）:")
    print(documents[0].page_content[:500] + "...")
